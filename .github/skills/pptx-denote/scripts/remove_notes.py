#!/usr/bin/env python3
"""Remove all slide notes from a .pptx file while keeping the slides intact.

Usage:
    python remove_notes.py "input.pptx" [output.pptx]

If no output path is given, "<input>-denoted.pptx" is written next to the input.
Requires: pip install python-pptx
"""

import sys
from pathlib import Path

from pptx import Presentation


def remove_notes(input_path: Path, output_path: Path) -> int:
    prs = Presentation(str(input_path))
    cleared = 0

    for slide in prs.slides:
        if not slide.has_notes_slide:
            continue
        text_frame = slide.notes_slide.notes_text_frame
        if text_frame.text:
            text_frame.clear()
            cleared += 1

    prs.save(str(output_path))
    return cleared


def main() -> None:
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if not input_path.is_file():
        print(f"Error: file not found: {input_path}")
        sys.exit(1)

    if len(sys.argv) >= 3:
        output_path = Path(sys.argv[2])
    else:
        output_path = input_path.with_name(f"{input_path.stem}-denoted{input_path.suffix}")

    cleared = remove_notes(input_path, output_path)
    print(f"Cleared notes on {cleared} slide(s).")
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
